"""
Extractor simple para .docx y .pptx. Crea archivos de texto en ./extracted/ con el contenido extraído.

Uso:
  python3 extract_text.py

Instalar dependencias:
  pip install python-docx python-pptx
"""
import os
from pathlib import Path

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

BASE = Path(__file__).parent
EXTRACT_DIR = BASE / 'extracted'
EXTRACT_DIR.mkdir(exist_ok=True)

def extract_docx(path: Path):
    if Document is None:
        print('python-docx no instalado. Ejecuta: pip install python-docx')
        return
    doc = Document(path)
    lines = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            lines.append(text)
    out = EXTRACT_DIR / (path.stem + '.txt')
    out.write_text('\n\n'.join(lines), encoding='utf-8')
    print('Wrote', out)

def extract_pptx(path: Path):
    if Presentation is None:
        print('python-pptx no instalado. Ejecuta: pip install python-pptx')
        return
    pres = Presentation(path)
    lines = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = shape.text.strip()
                if text:
                    lines.append(text)
    out = EXTRACT_DIR / (path.stem + '.txt')
    out.write_text('\n\n'.join(lines), encoding='utf-8')
    print('Wrote', out)

def main():
    workspace = BASE.parent
    candidates = list(workspace.glob('*.docx')) + list(workspace.glob('*.pptx'))
    if not candidates:
        print('No .docx or .pptx encontrados en', workspace)
        return
    for p in candidates:
        print('Processing', p.name)
        if p.suffix.lower() == '.docx':
            extract_docx(p)
        elif p.suffix.lower() == '.pptx':
            extract_pptx(p)

if __name__ == '__main__':
    main()
